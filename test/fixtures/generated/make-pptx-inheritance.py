#!/usr/bin/env python3
"""Generate the python-pptx fixture — the placeholder inheritance chain, unfaked.

`test/fixtures/corpus/make-corpus-pptx.py` writes `<a:buChar>` into every paragraph it
wants bulleted, so the corpus can only ever prove that slimdoc reads a bullet that is
written down. PowerPoint does not write bullets down: a content placeholder inherits its
glyph from the layout, which inherits from the master's `<p:bodyStyle>`. A deck whose
bullets are all explicit therefore agrees with an extractor that only reads explicit ones.

python-pptx ships PowerPoint's own default template, so the master, the layouts and the
placeholder `idx`/`type` wiring here are Microsoft's rather than ours. The deck contains
one of each case the resolution has to tell apart:

  slide 1  title placeholder             -> master titleStyle, no bullet
           body placeholder, plain       -> master bodyStyle, INHERITS a bullet
           body placeholder, buNone      -> explicitly off, must stay off
  slide 2  plain text box (no p:ph)      -> master otherStyle, no bullet
           body placeholder, explicit    -> explicit buChar, on
  slide 3  body placeholder, nested      -> lvl 0/1/2 all inheriting

The buNone and explicit-buChar cases are set through lxml, because python-pptx's API has
no bullet property — which is itself the point: those are the two cases a deck author
sets deliberately, and everything else inherits.

Deterministic: core property timestamps and every zip entry date are pinned, so
re-running produces a byte-identical file. Verify with:

    shasum -a 256 inherited-bullets.pptx && python3 make-pptx-inheritance.py && \\
    shasum -a 256 inherited-bullets.pptx

Usage:  python3 test/fixtures/generated/make-pptx-inheritance.py
"""

from __future__ import annotations

import datetime as dt
import os
import shutil
import zipfile

from pptx import Presentation
from pptx.util import Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "inherited-bullets.pptx")

# A fixed timestamp everywhere a clock would otherwise appear.
EPOCH = dt.datetime(2024, 1, 1, 0, 0, 0)
ZIP_DATE = (2024, 1, 1, 0, 0, 0)

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _pPr(paragraph):
    """The paragraph's `<a:pPr>`, created if the paragraph has none yet."""
    return paragraph._p.get_or_add_pPr()


def no_bullet(paragraph) -> None:
    """`<a:buNone/>` — the author turning inheritance off for this line."""
    _pPr(paragraph).append(paragraph._p.makeelement(f"{{{A_NS}}}buNone", {}))


def explicit_bullet(paragraph, char: str = "•") -> None:
    """`<a:buChar char="•"/>` — the author setting a glyph on this line."""
    _pPr(paragraph).append(
        paragraph._p.makeelement(f"{{{A_NS}}}buChar", {"char": char})
    )


def build(path: str) -> None:
    prs = Presentation()  # PowerPoint's default template: real master, real layouts
    title_and_content = prs.slide_layouts[1]
    blank = prs.slide_layouts[6]

    # -- slide 1: title, an inheriting body, and a body line turned off ------
    one = prs.slides.add_slide(title_and_content)
    one.shapes.title.text = "Propulsion"
    body = one.placeholders[1].text_frame

    body.text = "Warp field geometry"                   # inherits: expect a bullet
    body.add_paragraph().text = "Nacelle plasma injectors"   # inherits: expect a bullet

    off = body.add_paragraph()
    off.text = "This line opted out of the list"
    no_bullet(off)                                       # expect NO bullet

    # -- slide 2: a plain text box beside an explicitly bulleted body --------
    two = prs.slides.add_slide(blank)
    box = two.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(2))
    frame = box.text_frame
    frame.text = "A caption in a plain text box"         # otherStyle: expect NO bullet
    frame.add_paragraph().text = "A second caption line"  # otherStyle: expect NO bullet

    three = prs.slides.add_slide(title_and_content)
    three.shapes.title.text = "Refit sequence"
    seq = three.placeholders[1].text_frame
    seq.text = "Dock and power down"
    explicit_bullet(seq.paragraphs[0])                   # expect a bullet

    # -- slide 3's body also carries depth, all of it inheriting ------------
    for depth, text in ((1, "Drain the plasma manifold"), (2, "Cap the conduit")):
        para = seq.add_paragraph()
        para.text = text
        para.level = depth                               # inherits: expect bullets

    prs.core_properties.title = "Inherited bullets"
    prs.core_properties.author = "slimdoc fixtures"
    prs.core_properties.created = EPOCH
    prs.core_properties.modified = EPOCH
    prs.core_properties.last_modified_by = "slimdoc fixtures"
    prs.core_properties.revision = 1

    prs.save(path)
    _pin_zip_dates(path)


def _pin_zip_dates(path: str) -> None:
    """Rewrite every entry with a fixed date, so two runs are byte-identical.

    python-pptx stamps the clock into each zip entry. Sizes and order are preserved;
    only the date fields change, and the parts themselves are copied verbatim.
    """
    source = zipfile.ZipFile(path)
    temporary = path + ".tmp"
    with zipfile.ZipFile(temporary, "w", zipfile.ZIP_DEFLATED) as out:
        for item in source.infolist():
            pinned = zipfile.ZipInfo(item.filename, date_time=ZIP_DATE)
            pinned.compress_type = item.compress_type
            pinned.external_attr = 0o644 << 16
            out.writestr(pinned, source.read(item.filename))
    source.close()
    shutil.move(temporary, path)


if __name__ == "__main__":
    build(OUT)
    print(f"wrote {os.path.basename(OUT)} ({os.path.getsize(OUT)} bytes)")
